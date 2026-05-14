"""Generate a minimal sample .pptx template for demonstration purposes.

Run:
    python examples/create_sample_template.py
Produces:
    templates/quarterly_report_template.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def add_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "{{company}} — {{quarter}} Results"
    slide.placeholders[1].text = "Prepared by Finance"


def add_kpi_slide(prs: Presentation) -> None:
    """Slide with flat-placeholder table (KPI summary)."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    title.text_frame.text = "Key Metrics — {{quarter}}"
    title.text_frame.paragraphs[0].runs[0].font.size = Pt(24)

    rows, cols = 2, 3
    table = slide.shapes.add_table(
        rows, cols,
        Inches(0.5), Inches(1.0), Inches(9), Inches(1.5),
    ).table

    headers = ["Revenue", "Growth", "Gross Margin"]
    placeholders = ["{{revenue}}", "{{growth}}", "{{gross_margin}}"]

    for c, header in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = header
        cell.text_frame.paragraphs[0].runs[0].font.bold = True

    for c, ph in enumerate(placeholders):
        table.cell(1, c).text = ph


def add_regional_slide(prs: Presentation) -> None:
    """Slide 2: table filled by slide2_table1 key."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    title.text_frame.text = "Regional Breakdown"
    title.text_frame.paragraphs[0].runs[0].font.size = Pt(24)

    rows, cols = 6, 3  # 1 header + 5 data rows (will be filled from JSON)
    table = slide.shapes.add_table(
        rows, cols,
        Inches(0.5), Inches(1.0), Inches(9), Inches(3),
    ).table

    for c, header in enumerate(["Region", "Revenue", "YoY Growth"]):
        cell = table.cell(0, c)
        cell.text = header
        cell.text_frame.paragraphs[0].runs[0].font.bold = True


def add_metrics_slide(prs: Presentation) -> None:
    """Slide 3: table filled by slide3_table1 key."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    title.text_frame.text = "Customer Metrics"
    title.text_frame.paragraphs[0].runs[0].font.size = Pt(24)

    rows, cols = 4, 3
    table = slide.shapes.add_table(
        rows, cols,
        Inches(0.5), Inches(1.0), Inches(9), Inches(2.5),
    ).table

    for c, header in enumerate(["Metric", "Target", "Actual"]):
        cell = table.cell(0, c)
        cell.text = header
        cell.text_frame.paragraphs[0].runs[0].font.bold = True


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)
    add_kpi_slide(prs)
    add_regional_slide(prs)
    add_metrics_slide(prs)

    out = "templates/quarterly_report_template.pptx"
    prs.save(out)
    print(f"Template saved: {out}")


if __name__ == "__main__":
    main()
