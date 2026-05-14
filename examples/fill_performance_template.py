"""Fill Performance_By_Objective_template.pptx with new data.

Usage:
    python examples/fill_performance_template.py
    python examples/fill_performance_template.py --output output/my_report.pptx

Edit the DATA dict below to supply your numbers.  Wrap any value in
above_benchmark() to highlight it green (= above benchmark).
"""

from __future__ import annotations

import argparse
import copy
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor

TEMPLATE = Path("templates/Performance_By_Objective_template.pptx")
DEFAULT_OUTPUT = Path("output/Performance_By_Objective_filled.pptx")

GREEN = RGBColor(0x6A, 0xDF, 0x5B)


# ---------------------------------------------------------------------------
# Helper to mark a value as "above benchmark" (renders green + bold)
# ---------------------------------------------------------------------------

def above_benchmark(value: str) -> dict:
    return {"value": value, "above_benchmark": True}


# ---------------------------------------------------------------------------
# DATA – edit this section with your numbers
# ---------------------------------------------------------------------------

DATA = {
    # Date shown in both slide titles
    "date_range": "11/17-11/30",

    # ── Slide 1: Performance By Objective ──────────────────────────────────
    # Columns: Reach | Top Feed | Pulse Premiere
    # Rows (in order): Cost, Impressions, CPM, Reach, Clicks (destination),
    #                  CTR (destination), CPC (Destination), VTR, 6-Sec VTR,
    #                  Engagement Rate
    "slide1": {
        "reach": {
            "cost":              "$242,477.98",
            "impressions":       "29,435,586",
            "cpm":               "$8.24",
            "reach":             "6,956,446",
            "clicks":            "70,372",
            "ctr":               above_benchmark("0.24%"),
            "cpc":               "$3.45",
            "vtr":               above_benchmark("1.18%"),
            "vtr_6sec":          above_benchmark("2.09%"),
            "engagement_rate":   above_benchmark("0.39%"),
        },
        "top_feed": {
            "cost":              "$26,089.81",
            "impressions":       "4,899,506",
            "cpm":               above_benchmark("$5.32"),
            "reach":             "1,031,565",
            "clicks":            "12,606",
            "ctr":               above_benchmark("0.26%"),
            "cpc":               above_benchmark("$2.07"),
            "vtr":               "0.60%",
            "vtr_6sec":          "2.07%",
            "engagement_rate":   above_benchmark("0.37%"),
        },
        "pulse_premiere": {
            "cost":              "$99,256.56",
            "impressions":       "3,308,553",
            "cpm":               "$30.00",
            "reach":             "1,453,696",
            "clicks":            "7,314",
            "ctr":               "0.22%",
            "cpc":               "$13.57",
            "vtr":               "0.66%",
            "vtr_6sec":          "2.35%",
            "engagement_rate":   "0.47%",
        },
    },

    # ── Slide 2: Performance By Creative ───────────────────────────────────
    # Columns: Creative | Cost | Impressions | CPM | CTR | CPC | 6s VTR | VTR | ER
    "slide2": {
        "creatives": [
            {
                "name":        "GREEN_PRODUCT_LANDO_VIDEO_15S",
                "cost":        "$157,346.85",
                "impressions": "21,898,963",
                "cpm":         "$7.19",
                "ctr":         above_benchmark("0.24%"),
                "cpc":         "$2.99",
                "vtr_6sec":    above_benchmark("1.97%"),
                "vtr":         "0.54%",
                "er":          above_benchmark("0.38%"),
            },
            {
                "name":        "GREEN_PRODUCT_LANDO_VIDEO_6S",
                "cost":        "$111,220.94",
                "impressions": "12,436,129",
                "cpm":         "$8.94",
                "ctr":         above_benchmark("0.24%"),
                "cpc":         "$3.67",
                "vtr_6sec":    above_benchmark("2.29%"),
                "vtr":         above_benchmark("2.07%"),
                "er":          above_benchmark("0.41%"),
            },
            {
                "name":        "GREEN_LIFESTYLE_4WHEEL_LANDO_VIDEO_15S",
                "cost":        "$99,256.56",
                "impressions": "3,308,553",
                "cpm":         "$30.00",
                "ctr":         above_benchmark("0.22%"),
                "cpc":         "$13.57",
                "vtr_6sec":    above_benchmark("2.35%"),
                "vtr":         "0.66%",
                "er":          above_benchmark("0.47%"),
            },
        ]
    },
}


# ---------------------------------------------------------------------------
# Core cell-writing logic
# ---------------------------------------------------------------------------

def _write_cell(cell, entry) -> None:
    """Write a value into a table cell, preserving font size.

    *entry* is either a plain string or a dict from above_benchmark().
    """
    if isinstance(entry, dict):
        text = entry["value"]
        highlight = entry.get("above_benchmark", False)
    else:
        text = str(entry)
        highlight = False

    tf = cell.text_frame
    if not tf.paragraphs:
        return

    para = tf.paragraphs[0]

    # Collect all runs; keep the first one to preserve size, drop the rest
    runs = para.runs
    if not runs:
        run = para.add_run()
    else:
        run = runs[0]
        for extra in runs[1:]:
            # Clear extra runs by emptying their text
            extra.text = ""

    run.text = text

    from pptx.oxml.ns import qn

    if highlight:
        run.font.bold = True
        run.font.color.rgb = GREEN
    else:
        run.font.bold = None          # inherit from cell/table style
        # Remove any explicit solidFill so the cell inherits its theme color
        rPr = run._r.get_or_add_rPr()
        solidFill = rPr.find(qn("a:solidFill"))
        if solidFill is not None:
            rPr.remove(solidFill)


# ---------------------------------------------------------------------------
# Slide fillers
# ---------------------------------------------------------------------------

def _fill_slide1(slide, data: dict) -> None:
    date_range = data["date_range"]
    obj_data = data["slide1"]

    # Update date in title placeholder (para index 1, run 0)
    for shape in slide.shapes:
        if shape.has_text_frame:
            paras = shape.text_frame.paragraphs
            if len(paras) > 1 and paras[0].text.strip() == "Performance By Objective":
                paras[1].runs[0].text = date_range
                break

    # Find the table
    table = next(s.table for s in slide.shapes if s.has_table)

    # Row order matches the template (rows 1-10, skipping header row 0)
    columns = ["reach", "top_feed", "pulse_premiere"]
    metrics = [
        "cost", "impressions", "cpm", "reach", "clicks",
        "ctr", "cpc", "vtr", "vtr_6sec", "engagement_rate",
    ]

    for col_idx, col_key in enumerate(columns, start=1):
        col_data = obj_data[col_key]
        for row_idx, metric_key in enumerate(metrics, start=1):
            _write_cell(table.cell(row_idx, col_idx), col_data[metric_key])


def _fill_slide2(slide, data: dict) -> None:
    date_range = data["date_range"]
    creatives = data["slide2"]["creatives"]

    # Update date
    for shape in slide.shapes:
        if shape.has_text_frame:
            paras = shape.text_frame.paragraphs
            if len(paras) > 1 and paras[0].text.strip() == "Performance By Creative":
                if paras[1].runs:
                    paras[1].runs[0].text = date_range
                break

    table = next(s.table for s in slide.shapes if s.has_table)

    # Column order matches the template header row
    fields = ["name", "cost", "impressions", "cpm", "ctr", "cpc", "vtr_6sec", "vtr", "er"]

    for row_idx, creative in enumerate(creatives, start=1):
        if row_idx >= len(table.rows):
            break
        for col_idx, field in enumerate(fields):
            _write_cell(table.cell(row_idx, col_idx), creative[field])


# ---------------------------------------------------------------------------
# Main entry point
# ---------------------------------------------------------------------------

def generate(data: dict, output_path: Path) -> Path:
    prs = Presentation(str(TEMPLATE))

    _fill_slide1(prs.slides[0], data)
    _fill_slide2(prs.slides[1], data)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path.resolve()


def main() -> None:
    parser = argparse.ArgumentParser(description="Fill Performance By Objective template")
    parser.add_argument("-o", "--output", default=str(DEFAULT_OUTPUT),
                        help=f"Output path (default: {DEFAULT_OUTPUT})")
    args = parser.parse_args()

    result = generate(DATA, Path(args.output))
    print(f"Saved: {result}")


if __name__ == "__main__":
    main()
