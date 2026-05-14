"""Core slide generation logic: load a .pptx template, fill tables, save output."""

from __future__ import annotations

import copy
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Pt

from .table_filler import fill_table


class SlideGenerator:
    """Fill tables in a PowerPoint template with supplied data and write a new deck."""

    def __init__(self, template_path: str | Path) -> None:
        self.template_path = Path(template_path)
        if not self.template_path.exists():
            raise FileNotFoundError(f"Template not found: {self.template_path}")
        self._prs = Presentation(str(self.template_path))

    def generate(
        self,
        data: dict[str, Any],
        output_path: str | Path,
    ) -> Path:
        """Fill the template with *data* and write the result to *output_path*.

        Parameters
        ----------
        data:
            Mapping of placeholder keys to values.  Two schemas are supported:

            1. **Flat key→value** – replaces every ``{{KEY}}`` placeholder found
               in any table cell or text frame.

               Example::

                   {"revenue": "1,234", "growth": "12%"}

            2. **Table-indexed** – targets a specific table by its 1-based index
               on a slide.  The key format is ``slide<N>_table<M>`` and the value
               is a 2-D list (list of rows, each row a list of cell values).

               Example::

                   {"slide1_table1": [["Q1", "Q2"], ["100", "200"]]}

        output_path:
            Destination ``.pptx`` file path.  Parent directories are created
            automatically.

        Returns
        -------
        Path
            Resolved path to the written file.
        """
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)

        prs = copy.deepcopy(self._prs)

        table_data = {k: v for k, v in data.items() if isinstance(v, list)}
        flat_data = {k: str(v) for k, v in data.items() if not isinstance(v, list)}

        for slide_idx, slide in enumerate(prs.slides, start=1):
            table_counter = 0
            for shape in slide.shapes:
                if shape.has_table:
                    table_counter += 1
                    key = f"slide{slide_idx}_table{table_counter}"
                    if key in table_data:
                        fill_table(shape.table, table_data[key])
                    elif flat_data:
                        fill_table(shape.table, None, placeholders=flat_data)

                if shape.has_text_frame and flat_data:
                    _replace_in_text_frame(shape.text_frame, flat_data)

        prs.save(str(output_path))
        return output_path.resolve()


def _replace_in_text_frame(text_frame: Any, replacements: dict[str, str]) -> None:
    for para in text_frame.paragraphs:
        for run in para.runs:
            for key, value in replacements.items():
                placeholder = f"{{{{{key}}}}}"
                if placeholder in run.text:
                    run.text = run.text.replace(placeholder, value)
