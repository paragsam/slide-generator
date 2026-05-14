"""Utilities for writing values into pptx Table objects."""

from __future__ import annotations

from typing import Any


def fill_table(
    table: Any,
    rows: list[list[str]] | None,
    *,
    placeholders: dict[str, str] | None = None,
) -> None:
    """Write values into *table*.

    Two modes:
    - ``rows`` provided: write the 2-D list directly into cells (row-major,
      skips header row if the table has one more row than the data).
    - ``placeholders`` provided: replace ``{{KEY}}`` tokens in every cell.
    """
    if rows is not None:
        _fill_by_rows(table, rows)
    elif placeholders:
        _fill_by_placeholders(table, placeholders)


def _fill_by_rows(table: Any, rows: list[list[str]]) -> None:
    table_row_count = len(table.rows)
    data_row_count = len(rows)

    # If template has a header row, start writing data from row index 1
    start_row = 1 if table_row_count == data_row_count + 1 else 0

    for r_idx, row_values in enumerate(rows):
        tbl_row_idx = start_row + r_idx
        if tbl_row_idx >= table_row_count:
            break
        for c_idx, value in enumerate(row_values):
            if c_idx >= len(table.columns):
                break
            _set_cell_text(table.cell(tbl_row_idx, c_idx), str(value))


def _fill_by_placeholders(table: Any, placeholders: dict[str, str]) -> None:
    for row in table.rows:
        for cell in row.cells:
            for key, value in placeholders.items():
                token = f"{{{{{key}}}}}"
                if token in cell.text:
                    _replace_cell_text(cell, token, value)


def _set_cell_text(cell: Any, text: str) -> None:
    tf = cell.text_frame
    if not tf.paragraphs:
        return
    para = tf.paragraphs[0]
    if para.runs:
        run = para.runs[0]
        # Preserve font size / bold if already set
        run.text = text
        for extra_run in para.runs[1:]:
            extra_run.text = ""
    else:
        from pptx.util import Pt
        run = para.add_run()
        run.text = text


def _replace_cell_text(cell: Any, token: str, value: str) -> None:
    for para in cell.text_frame.paragraphs:
        for run in para.runs:
            if token in run.text:
                run.text = run.text.replace(token, value)
