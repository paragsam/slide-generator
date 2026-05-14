"""Tests for SlideGenerator."""

import json
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from slide_generator import SlideGenerator


def _first_table(slide):
    for shape in slide.shapes:
        if shape.has_table:
            return shape.table
    raise AssertionError("No table found on slide")


FIXTURES = Path(__file__).parent / "fixtures"


def _make_template(tmp_path: Path, rows: int = 3, cols: int = 2) -> Path:
    """Create a minimal single-slide .pptx template with one table."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    table = slide.shapes.add_table(
        rows, cols, Inches(1), Inches(1), Inches(8), Inches(3)
    ).table

    # Header row
    for c in range(cols):
        table.cell(0, c).text = f"Header {c + 1}"

    # Data rows with placeholders
    for r in range(1, rows):
        for c in range(cols):
            table.cell(r, c).text = f"{{{{r{r}c{c + 1}}}}}"

    out = tmp_path / "template.pptx"
    prs.save(str(out))
    return out


class TestFlatPlaceholders:
    def test_replaces_placeholder_in_table(self, tmp_path):
        template = _make_template(tmp_path, rows=2, cols=2)
        gen = SlideGenerator(template)
        output = gen.generate({"r1c1": "Alpha", "r1c2": "Beta"}, tmp_path / "out.pptx")

        prs = Presentation(str(output))
        table = _first_table(prs.slides[0])
        assert table.cell(1, 0).text == "Alpha"
        assert table.cell(1, 1).text == "Beta"


class TestTableIndexedData:
    def test_fills_table_by_key(self, tmp_path):
        template = _make_template(tmp_path, rows=3, cols=2)
        data = {"slide1_table1": [["10", "20"], ["30", "40"]]}
        gen = SlideGenerator(template)
        output = gen.generate(data, tmp_path / "out.pptx")

        prs = Presentation(str(output))
        table = _first_table(prs.slides[0])
        assert table.cell(1, 0).text == "10"
        assert table.cell(2, 1).text == "40"

    def test_unknown_key_leaves_template_unchanged(self, tmp_path):
        template = _make_template(tmp_path, rows=2, cols=1)
        gen = SlideGenerator(template)
        output = gen.generate({"slide9_table9": [["X"]]}, tmp_path / "out.pptx")

        prs = Presentation(str(output))
        table = _first_table(prs.slides[0])
        assert table.cell(1, 0).text == "{{r1c1}}"


class TestOutputPath:
    def test_creates_parent_dirs(self, tmp_path):
        template = _make_template(tmp_path, rows=2, cols=1)
        gen = SlideGenerator(template)
        nested = tmp_path / "a" / "b" / "out.pptx"
        result = gen.generate({}, nested)
        assert result.exists()

    def test_missing_template_raises(self, tmp_path):
        with pytest.raises(FileNotFoundError):
            SlideGenerator(tmp_path / "nonexistent.pptx")
